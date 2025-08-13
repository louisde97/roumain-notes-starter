from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from sqlmodel import SQLModel, Field, Session, create_engine, select
from typing import Optional, List
import os
import io
from datetime import date, timedelta
from docx import Document as DocxDocument
from pptx import Presentation

# ---------- Config ----------
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./local.db")
engine = create_engine(DATABASE_URL, echo=False)

CORS_ORIGINS = os.getenv("CORS_ORIGINS", "*").split(",")

app = FastAPI(title="Roumain Notes — Backend (MVP)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[o.strip() for o in CORS_ORIGINS if o.strip()],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------- Models ----------
class Lesson(SQLModel, table=True):
    id: Optional[int] = Field(default=None, primary_key=True)
    title: Optional[str] = None

class Concept(SQLModel, table=True):
    id: Optional[int] = Field(default=None, primary_key=True)
    type: str  # vocab | conjugaison | prononciation | grammaire | expression | orthographe | culture
    key_lemma: str
    title: str
    data: Optional[str] = None  # JSON text for MVP simplicity
    __table_args__ = ({"sqlite_autoincrement": True},)

class Occurrence(SQLModel, table=True):
    id: Optional[int] = Field(default=None, primary_key=True)
    concept_id: int = Field(foreign_key="concept.id")
    lesson_id: int = Field(foreign_key="lesson.id")
    is_revision: bool = False
    extract: Optional[str] = None

class Card(SQLModel, table=True):
    id: Optional[int] = Field(default=None, primary_key=True)
    concept_id: int = Field(foreign_key="concept.id")
    front: str
    back: str
    next_review_date: date = Field(default_factory=date.today)
    interval_days: int = 0
    ease_factor: float = 2.5
    repetitions: int = 0

def init_db():
    SQLModel.metadata.create_all(engine)

@app.on_event("startup")
def on_startup():
    init_db()

# ---------- Heuristics ----------
GRAMMAR_HINTS = [
    "genitiv", "dativ", "acuzativ", "nominativ",
    "articol", "plural", "singular", "masculin", "feminin", "neutru",
    "prezent", "imperfect", "perfect compus", "viitor", "conjunctiv", "imperativ",
    "pronume", "substantiv", "adjectiv", "adverb", "prepozi", "conjunc", "interjec"
]

def is_conjugation_block(lines: List[str]) -> bool:
    persons = ["eu", "tu", "el", "ea", "noi", "voi", "ei", "ele"]
    count = sum(any(l.strip().lower().startswith(p + " ") for p in persons) for l in lines)
    return count >= 3

def parse_vocab_line(line: str):
    # Common separators: " – ", " - ", " — ", " : ", " = "
    for sep in [" – ", " — ", " - ", " : ", " = "]:
        if sep in line:
            left, right = line.split(sep, 1)
            ro = left.strip().strip("•-*–—:").strip()
            fr = right.strip()
            if ro and fr:
                return ro, fr
    return None

def categorize_line(line: str) -> Optional[str]:
    low = line.lower()
    if any(k in low for k in ["pronun", "prononci", "accent"]):
        return "prononciation"
    if any(k in low for k in GRAMMAR_HINTS):
        return "grammaire"
    if parse_vocab_line(line):
        return "vocab"
    if any(k in low for k in ["expres", "locuțiune", "idiom"]):
        return "expression"
    if any(k in low for k in ["ortograf", "diacrit", "scriere"]):
        return "orthographe"
    return None

def extract_from_docx(content: bytes) -> List[str]:
    f = io.BytesIO(content)
    doc = DocxDocument(f)
    lines = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            lines.append(t)
    # tables too
    for table in doc.tables:
        for row in table.rows:
            row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text)
            if row_text.strip():
                lines.append(row_text.strip())
    return lines

def extract_from_pptx(content: bytes) -> List[str]:
    f = io.BytesIO(content)
    prs = Presentation(f)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    for line in t.splitlines():
                        if line.strip():
                            lines.append(line.strip())
    return lines

def upsert_concept(session: Session, ctype: str, key_lemma: str, title: str, data: Optional[str]) -> Concept:
    stmt = select(Concept).where(Concept.type == ctype, Concept.key_lemma == key_lemma)
    existing = session.exec(stmt).first()
    if existing:
        return existing
    obj = Concept(type=ctype, key_lemma=key_lemma, title=title, data=data)
    session.add(obj)
    session.commit()
    session.refresh(obj)
    return obj

# ---------- API ----------
@app.get("/health")
def health():
    return {"ok": True}

@app.post("/api/upload")
async def upload(file: UploadFile = File(...)):
    name = (file.filename or "lecon").rsplit(".", 1)[0]
    content = await file.read()
    ext = (file.filename or "").lower().rsplit(".", 1)[-1]
    if ext == "docx":
        lines = extract_from_docx(content)
    elif ext == "pptx":
        lines = extract_from_pptx(content)
    else:
        return {"error": "Format non supporté. Utilise .docx ou .pptx"}

    # Insert lesson
    with Session(engine) as session:
        lesson = Lesson(title=name)
        session.add(lesson)
        session.commit()
        session.refresh(lesson)

        # Pass 1: blocks by empties or headings (simple)
        blocks: List[List[str]] = []
        current: List[str] = []
        for ln in lines:
            if not ln.strip():
                if current:
                    blocks.append(current)
                    current = []
                continue
            current.append(ln)
        if current:
            blocks.append(current)

        extracted = []
        for blk in blocks:
            # conjugation block?
            if is_conjugation_block(blk):
                verb_name = None
                for line in blk[:3]:
                    # try find infinitive mention like "a merge", else take token after "eu"
                    if " a " in line.lower() or line.strip().lower().startswith("a "):
                        verb_name = line.strip()
                        break
                if not verb_name:
                    verb_name = "Conjugare"
                conj = {}
                for l in blk:
                    parts = l.split(":", 1)
                    if len(parts) == 2:
                        conj[parts[0].strip().lower()] = parts[1].strip()
                    else:
                        # try space split
                        sp = l.split(maxsplit=1)
                        if len(sp) == 2 and sp[0].lower() in ["eu","tu","el","ea","noi","voi","ei","ele"]:
                            conj[sp[0].lower()] = sp[1].strip()
                data = {"type":"conjugaison","verbe":verb_name,"paradigme":conj}
                c = upsert_concept(session, "conjugaison", key_lemma=verb_name.lower(), title=f"Conjugaison · {verb_name}", data=json.dumps(data, ensure_ascii=False))
                occ = Occurrence(concept_id=c.id, lesson_id=lesson.id, is_revision=True, extract="\n".join(blk)) if c else None
                if occ:
                    session.add(occ)
                    extracted.append({"type":"conjugaison","title":c.title,"is_revision":True})
                continue

            # lines inside block
            for line in blk:
                cat = categorize_line(line) or "culture"
                if cat == "vocab":
                    pair = parse_vocab_line(line)
                    if pair:
                        ro, fr = pair
                        key = ro.lower()
                        data = {"type":"vocab","ro":ro,"fr":fr}
                        c = upsert_concept(session, "vocab", key_lemma=key, title=f"{ro} — {fr}", data=json.dumps(data, ensure_ascii=False))
                        # is revision if already existed before upsert (simplified heuristic handled above)
                        # We cannot know here, so re-check count of occurrences:
                        prev_occ = session.exec(select(Occurrence).where(Occurrence.concept_id == c.id)).first()
                        is_rev = prev_occ is not None
                        occ = Occurrence(concept_id=c.id, lesson_id=lesson.id, is_revision=is_rev, extract=line)
                        session.add(occ)
                        extracted.append({"type":"vocab","title":c.title,"is_revision":is_rev})
                elif cat in ["prononciation","grammaire","expression","orthographe","culture"]:
                    key = line.lower()[:80]
                    data = {"type":cat,"note":line}
                    c = upsert_concept(session, cat, key_lemma=key, title=line[:80], data=json.dumps(data, ensure_ascii=False))
                    prev_occ = session.exec(select(Occurrence).where(Occurrence.concept_id == c.id)).first()
                    is_rev = prev_occ is not None
                    occ = Occurrence(concept_id=c.id, lesson_id=lesson.id, is_revision=is_rev, extract=line)
                    session.add(occ)
                    extracted.append({"type":cat,"title":c.title,"is_revision":is_rev})

        session.commit()

        # group by category for return
        by_cat = {}
        for item in extracted:
            by_cat.setdefault(item["type"], []).append({"title": item["title"], "is_revision": item["is_revision"]})

        return {"lesson_id": lesson.id, "title": lesson.title, "notions": by_cat}

@app.get("/api/lessons")
def list_lessons():
    with Session(engine) as session:
        lessons = session.exec(select(Lesson)).all()
        return [{"id": l.id, "title": l.title} for l in lessons]

@app.get("/api/lessons/{lesson_id}/notions")
def get_lesson_notions(lesson_id: int):
    with Session(engine) as session:
        stmt = select(Occurrence, Concept).where(Occurrence.lesson_id == lesson_id).join(Concept, Concept.id == Occurrence.concept_id)
        rows = session.exec(stmt).all()
        by_cat = {}
        for occ, c in rows:
            by_cat.setdefault(c.type, []).append({"title": c.title, "is_revision": occ.is_revision})
        return {"lesson_id": lesson_id, "notions": by_cat}

# Simple SM-2 endpoint (placeholder)
@app.post("/api/review/{concept_id}/{quality}")
def review(concept_id: int, quality: int):
    # quality: 0..5
    with Session(engine) as session:
        card = session.exec(select(Card).where(Card.concept_id == concept_id)).first()
        if not card:
            card = Card(concept_id=concept_id, front=\"\", back=\"\")
        # SM-2 logic (simplified)
        q = max(0, min(5, int(quality)))
        if q < 3:
            card.repetitions = 0
            card.interval_days = 1
        else:
            if card.repetitions == 0:
                card.interval_days = 1
            elif card.repetitions == 1:
                card.interval_days = 6
            else:
                card.interval_days = int(card.interval_days * card.ease_factor)
            card.ease_factor = max(1.3, card.ease_factor + (0.1 - (5 - q) * (0.08 + (5 - q) * 0.02)))
            card.repetitions += 1
        card.next_review_date = date.today() + timedelta(days=card.interval_days)
        session.add(card)
        session.commit()
        session.refresh(card)
        return {"ok": True, "next_review_date": str(card.next_review_date), "ease": card.ease_factor, "reps": card.repetitions, "interval_days": card.interval_days}
